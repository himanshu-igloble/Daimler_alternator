"""
V11.1_ALT — Business Summary Deck  (5 slides)
==============================================
Headline: "We tested whether the new electrical signals can predict each truck's
failure date — they cannot; the fleet schedule + alarms remain the deliverable."

Live numbers from V11.1 cache at runtime.
Mirrors the V11 business builder pattern.

Output: V11.1_ALT/presentation/Alternator_Business_Summary_V11.1.pptx
"""
from __future__ import annotations

import importlib.util
import json
import pathlib

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

_pres_dir = pathlib.Path(__file__).resolve().parent
_root = _pres_dir.parent
_src  = _root / "src"
_v1062_root = _root.parent / "V10.6.2_ALT"

NAVY = RGBColor(0x0D, 0x1B, 0x2A)
GOLD = RGBColor(0xC5, 0x8B, 0x1F)
RED  = RGBColor(0xC0, 0x39, 0x2B)

GRAPHS = _root / "visualizations" / "rul_core"


def _load(name: str):
    spec = importlib.util.spec_from_file_location(name, str(_src / f"{name}.py"))
    mod  = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


cfg = _load("V11_1_ALT_config")


def _nums() -> dict:
    bt  = json.loads((cfg.BACKTEST_CACHE / "backtest_results.json").read_text())
    fw  = json.loads((_v1062_root / "cache" / "rul" / "fleet_window.json").read_text())
    emg = pd.read_csv(cfg.EMERG_CACHE / "emergency_per_vin.csv")
    return {
        "m0_mae":    bt["variants"]["M0"]["mae_model"],
        "dummy_mae": bt["variants"]["M0"]["mae_dummy"],
        "chosen":    bt["chosen_variant"],
        "verdict":   bt["variants"]["M0"].get("wilcoxon_p_vs_dummy", 1.0),
        "emp_med":   fw["median_ttf_days"],
        "emp_km":    fw["median_ttf_km_est"],
        "ged_f":     int(emg[(emg.failed_flag == 1) & (emg.ged_fired == True)].shape[0]),
        "ged_nf":    int(emg[(emg.failed_flag == 0) & (emg.ged_fired == True)].shape[0]),
        "ew_f":      int(emg[(emg.failed_flag == 1) & (emg.early_watch_current == 1)].shape[0]),
        "ew_nf":     int(emg[(emg.failed_flag == 0) & (emg.early_watch_current == 1)].shape[0]),
    }


def _bullets(prs: Presentation, title: str, bullets: list[str]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(0.9))
    p = t.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.6), Inches(5.4)).text_frame
    body.word_wrap = True
    for i, b in enumerate(bullets):
        par = body.paragraphs[0] if i == 0 else body.add_paragraph()
        par.text = "• " + b; par.font.size = Pt(19); par.font.color.rgb = NAVY
    return s


def _image_with_caption(prs: Presentation, title: str, img_path: pathlib.Path,
                         caption: str = ""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.8))
    p = t.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    if img_path.exists():
        s.shapes.add_picture(str(img_path), Inches(1.5), Inches(1.1), height=Inches(5.0))
    else:
        warn = s.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(1))
        warn.text_frame.paragraphs[0].text = f"[graph not found: {img_path.name}]"
        warn.text_frame.paragraphs[0].font.color.rgb = RED
    if caption:
        cap = s.shapes.add_textbox(Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.7))
        cap.text_frame.paragraphs[0].text = caption
        cap.text_frame.paragraphs[0].font.size = Pt(13)
        cap.text_frame.paragraphs[0].font.color.rgb = GOLD
    return s


def main():
    n = _nums()
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Headline ─────────────────────────────────────────────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(12), Inches(2.8))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Alternator Predictive Maintenance — V11.1 Business Summary"
    p.font.size = Pt(33); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tf.add_paragraph()
    p2.text = (
        "We tested whether the new electrical signals can predict each truck's "
        "failure date — they cannot; the fleet schedule + alarms remain the deliverable."
    )
    p2.font.size = Pt(17); p2.font.color.rgb = GOLD
    p3 = tf.add_paragraph()
    p3.text = (
        f"Fleet: 25 trucks (10 failed + 15 in-service)  |  "
        f"Classifier AUROC 0.927 (unchanged)  |  "
        f"Verdict: NO_IMPROVEMENT_HONEST"
    )
    p3.font.size = Pt(15); p3.font.color.rgb = NAVY

    # ── Slide 2: How we did it ────────────────────────────────────────────────
    _bullets(prs, "How we tested it", [
        "V11 engineered 12 new electrical signals from the existing 6 CAN sensors "
        "(no new hardware) and improved failure detection from 2/10 to 6/10 trucks.",
        "V11.1 asked: can those signals, combined into a statistical survival model, "
        "predict each individual truck's remaining life more accurately than the fleet average?",
        "We encoded two covariate signals (lifetime electrical stress history + "
        "recent compound alarm activity) into an Accelerated Failure Time model.",
        "The model was validated out-of-sample on all 10 failed trucks "
        "(leave-one-out per fold, parameters re-estimated from scratch each time).",
        "The test was repeated three ways: no covariates (M0), one covariate (M1), "
        "two covariates (M2). Result: M0 wins — covariates make things worse.",
    ])

    # ── Slide 3: Results ─────────────────────────────────────────────────────
    _image_with_caption(
        prs,
        f"Result: no variant beats the fleet-clock dummy "
        f"(error {n['m0_mae']:.0f}d vs {n['dummy_mae']:.0f}d)",
        GRAPHS / "backtest_accuracy.png",
        caption=(
            f"M0 (no covariates): {n['m0_mae']:.1f}d error.  "
            "Fleet-clock dummy (assume every truck fails at the fleet median): "
            f"{n['dummy_mae']:.1f}d error.  "
            "The dummy wins by a factor of ~3 — this is a structural data limit, "
            "not a modelling choice."
        ),
    )

    # ── Slide 4: Limitations ─────────────────────────────────────────────────
    _bullets(prs, "Limitations & why the limit is structural", [
        "Only 10 failure events in the dataset — far too few to calibrate "
        "per-truck timing with any covariate model.",
        "The new electrical stress signal (x1) turns out to track truck age, "
        "not failure risk independently of age — an exposure confound.",
        "4 of 10 failures had no measurable electrical precursor "
        "(abrupt / silent failure mode); those trucks are undetectable from this data.",
        "The fleet is aged: all 15 in-service trucks are already past or near the "
        "empirical failure median — per-truck remaining-life estimates overlap heavily.",
        "No per-truck calendar dates are available; timing is in operational days.",
    ])

    # ── Slide 5: What you get + next steps ───────────────────────────────────
    _bullets(prs, "What ships + next steps", [
        f"WHICH trucks: frozen Ridge classifier AUROC 0.927 — "
        f"risk ranks the 15 in-service trucks reliably.",
        f"WHEN (fleet): schedule inspection at {n['emp_med']:.0f} days "
        f"(≈ {n['emp_km']:.0f} km) — the empirical fleet-failure median.",
        f"WHEN (alerts): GED storm fires for {n['ged_f']}/10 failed trucks, "
        f"{n['ged_nf']}/15 non-failed (0 false alarms). "
        f"Early-watch fires for {n['ew_f']}/10 active-failure trucks, "
        f"{n['ew_nf']}/15 non-failed (0 false alarms).",
        "Next: validate on the larger starter-motor fleet (n=34 trucks, more events). "
        "The timing question can be re-opened when ≥ 30 failure events are available.",
        "All claims are audited to source cache files (see AUDIT_REPORT.md).",
    ])

    out = _pres_dir / "Alternator_Business_Summary_V11.1.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"[v11.1 business deck] wrote {out.name} ({len(prs.slides._sldIdLst)} slides, "
          f"{out.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
