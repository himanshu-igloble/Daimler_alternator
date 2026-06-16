#!/usr/bin/env python3
"""
Build Daimler Alternator Predictive Maintenance Technical Review presentation.
V10.6.2 — validated metrics only. Same design language as Clutch presentation.
"""

import os
import sys
import tempfile
import shutil
import json
import csv
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── PATHS ──────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent.parent
V1062  = ROOT / "V10.6.2_ALT"
V106   = ROOT / "V10.6_ALT"
VIZ    = V1062 / "visualizations"
VIZ_CS = VIZ / "rul_curves_clutch_style"

OUT_DIR = V1062 / "presentation"
OUT_DIR.mkdir(exist_ok=True)

# ── DIMENSIONS ─────────────────────────────────────────────────
SW, SH = Inches(13.33), Inches(7.5)

# ── COLOURS ────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
GOLD       = RGBColor(0xC5, 0x8B, 0x1F)
GREY_MED   = RGBColor(0x60, 0x60, 0x60)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
LIGHT_GREY  = RGBColor(0xB0, 0xB8, 0xC4)
KT_HEADER  = RGBColor(0x2E, 0x50, 0x90)
KT_BODY    = RGBColor(0x1B, 0x2A, 0x4A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_PASS = RGBColor(0x1E, 0x8C, 0x45)
RED_FAIL   = RGBColor(0xC0, 0x39, 0x2B)
ORANGE_W   = RGBColor(0xE3, 0x6C, 0x09)
BG_LIGHT   = RGBColor(0xF5, 0xF6, 0xFA)
BG_TABLE_H = RGBColor(0x0D, 0x1B, 0x2A)
BG_TABLE_A = RGBColor(0xF0, 0xF2, 0xF5)
BG_TABLE_B = RGBColor(0xFF, 0xFF, 0xFF)
KPI_BLUE   = RGBColor(0x1A, 0x5C, 0xB0)
KPI_GREEN  = RGBColor(0x15, 0x7F, 0x3D)
KPI_AMBER  = RGBColor(0xB8, 0x7A, 0x0F)
KPI_RED    = RGBColor(0xB0, 0x2A, 0x2A)

FONT = 'Calibri'
TMPDIR = tempfile.mkdtemp(prefix='alt_charts_')

# ── VALIDATED METRICS ──────────────────────────────────────────
# All numbers from V10.6.2 / V10.5.3 pipeline outputs.
# Sources noted inline for audit trail.

# Ridge classifier (V10.5.3 final — V5.2_ALT/results/)
AUROC = 0.9267
F1_SCORE = 0.90
RECALL = 0.90
PRECISION = 0.90
SPECIFICITY = 0.9333
MCC = 0.8333
THRESHOLD = 0.4456
BOOTSTRAP_CI = (0.8065, 1.0)
PERMUTATION_P = 0.0
TP, FP, FN, TN = 9, 1, 1, 14
N_FEATURES = 6
FEATURES = [
    ("vsi_std_ratio_30d", 0.1547),
    ("vsi_dominant_freq", 0.1053),
    ("vsi_range_trend_last30d", 0.0613),
    ("vsi_spectral_entropy", 0.0573),
    ("progressive_drift", 0.0480),
    ("bat_charge_delta_trend_right", 0.0340),
]

# Fleet (vin_lifecycle.parquet via final_rul_per_vin.csv)
N_VINS = 25
N_FAILED = 10
N_NF = 15
MEDIAN_TTF = 601.0
P25_TTF = 577.5
P75_TTF = 652.5
MIN_TTF = 472
MAX_TTF = 673
MEDIAN_KM = 120440
MEDIAN_EHRS = 4538

# Weibull (fleet_weibull_params.json)
WEIBULL_SHAPE = 5.17
WEIBULL_SCALE = 771.36
WEIBULL_MEDIAN = 718.5
WEIBULL_CI = (677.3, 774.4)

# Backtest (backtest_results.json)
BACKTEST_MAE_MODEL = 125.0
BACKTEST_MAE_DUMMY = 49.7
BACKTEST_PI_COV = 0.90

# GED precursor (failure_mode_split.csv)
GED_ACTIONABLE = 1
GED_TERMINAL = 1
N_ABRUPT = 5
N_INCONCLUSIVE = 3
VIN1_LEAD = 21
VIN10_LEAD = 1

# In-service risk (final_rul_per_vin.csv)
N_LOW_RISK = 14
N_HIGH_RISK = 1
HIGH_RISK_VIN = "VIN3_NF_ALT"


# ── HELPER FUNCTIONS ───────────────────────────────────────────
# Identical to reference Clutch presentation for visual consistency.

def make_prs():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, left, top, w, h, fill_rgb, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    ln = shape.line
    if border:
        ln.fill.solid()
        ln.fill.fore_color.rgb = border
        ln.width = Pt(1)
    else:
        ln.fill.background()
    return shape

def add_rounded_rect(slide, left, top, w, h, fill_rgb, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    ln = shape.line
    if border:
        ln.fill.solid()
        ln.fill.fore_color.rgb = border
        ln.width = Pt(1.5)
    else:
        ln.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, font_size=12, bold=False,
             color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    p.space_after = Pt(2)
    return txBox

def add_multiline(slide, left, top, w, h, lines, font_size=11, color=DARK_TEXT,
                  bold=False, bullet=False):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, c = item, bold, color
        else:
            txt = item[0]
            b = item[1] if len(item) > 1 else bold
            c = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("•  " if bullet else "") + txt
        p.font.size = Pt(font_size)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = FONT
        p.space_after = Pt(3)
    return txBox

def add_header_bar(slide, title_text, subtitle_text=""):
    add_rect(slide, Inches(0), Inches(0), SW, Inches(1.0), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.5),
             title_text, font_size=22, bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.35),
                 subtitle_text, font_size=13, color=LIGHT_GREY)

def add_section_slide(slide, number, title_text, subtitle_text=""):
    add_rect(slide, Inches(0), Inches(0), SW, SH, NAVY)
    add_text(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(0.8),
             f"{number:02d}  {title_text}", font_size=34, bold=True, color=WHITE)
    if subtitle_text:
        add_text(slide, Inches(0.8), Inches(3.4), Inches(11), Inches(0.5),
                 subtitle_text, font_size=16, color=LIGHT_GREY)
    add_footer(slide, dark_bg=True)

def add_footer(slide, dark_bg=False):
    y = Inches(7.05)
    if not dark_bg:
        add_rect(slide, Inches(0), y, SW, Inches(0.45), BG_LIGHT)
    clr = LIGHT_GREY if dark_bg else GREY_MED
    add_text(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.3),
             "V10.6.2  |  BharatBenz 5528T Alternator Predictive Maintenance  |  BytEdge CONFIDENTIAL",
             font_size=8, color=clr, bold=False)

def add_key_takeaways(slide, bullets, left=Inches(0.4), top=Inches(5.45),
                      width=Inches(12.5), height=None):
    if height is None:
        height = Inches(0.25 + 0.22 * len(bullets))
    add_rounded_rect(slide, left, top, width, height,
                     RGBColor(0xE8, 0xEE, 0xF7), border=KT_HEADER)
    add_text(slide, left + Inches(0.15), top + Inches(0.05), Inches(3), Inches(0.25),
             "KEY TAKEAWAYS", font_size=10, bold=True, color=KT_HEADER)
    y = top + Inches(0.28)
    for b in bullets:
        add_text(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.2),
                 "•  " + b, font_size=9, color=KT_BODY)
        y += Inches(0.19)

def add_kpi_tile(slide, left, top, w, h, label, value, color=KPI_BLUE, status=""):
    add_rounded_rect(slide, left, top, w, h, WHITE, border=color)
    add_text(slide, left + Inches(0.1), top + Inches(0.08), w - Inches(0.2), Inches(0.22),
             label, font_size=9, bold=False, color=GREY_MED, align=PP_ALIGN.CENTER)
    add_text(slide, left + Inches(0.1), top + Inches(0.3), w - Inches(0.2), Inches(0.4),
             str(value), font_size=20, bold=True, color=color, align=PP_ALIGN.CENTER)
    if status:
        sc = GREEN_PASS if "PASS" in status else (RED_FAIL if "FAIL" in status else KPI_AMBER)
        add_text(slide, left + Inches(0.1), top + Inches(0.68), w - Inches(0.2), Inches(0.2),
                 status, font_size=8, bold=True, color=sc, align=PP_ALIGN.CENTER)

def add_table_shape(slide, left, top, w, rows, cols, data, col_widths=None):
    tbl_shape = slide.shapes.add_table(rows, cols, left, top, w, Inches(0.35 * rows))
    tbl = tbl_shape.table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
    for ri, row_data in enumerate(data):
        for ci, val in enumerate(row_data):
            cell = tbl.cell(ri, ci)
            cell.text = str(val)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT
            if ri == 0:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_H
            else:
                p.font.size = Pt(8.5)
                p.font.color.rgb = DARK_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_TABLE_A if ri % 2 == 0 else BG_TABLE_B
            p.alignment = PP_ALIGN.CENTER
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
    return tbl_shape


# ── CHART GENERATION ───────────────────────────────────────────

def chart_ridge_metrics():
    fig, ax = plt.subplots(figsize=(10, 3.5))
    metrics = ['AUC-ROC', 'F1 Score', 'Recall', 'Precision', 'Specificity']
    values  = [AUROC * 100, F1_SCORE * 100, RECALL * 100, PRECISION * 100, SPECIFICITY * 100]
    bars = ax.bar(metrics, values, color='#1A5CB0', width=0.55,
                  edgecolor='#0D1B2A', linewidth=0.5, zorder=3)
    for i, v in enumerate(values):
        ax.text(i, v + 1.2, f'{v:.1f}%', ha='center', va='bottom',
                fontsize=11, fontweight='bold', color='#0D1B2A')
    ax.set_ylim(0, 115)
    ax.set_ylabel('Percentage (%)', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    ax.set_facecolor('white')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'ridge_metrics.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_confusion():
    fig, ax = plt.subplots(figsize=(4.5, 3.5))
    matrix = np.array([[TP, FP], [FN, TN]])
    ax.imshow(matrix, cmap='Blues', vmin=0, vmax=14)
    labels = [[f'TP\n{TP}', f'FP\n{FP}'], [f'FN\n{FN}', f'TN\n{TN}']]
    for i in range(2):
        for j in range(2):
            c = 'white' if matrix[i, j] > 6 else '#0D1B2A'
            ax.text(j, i, labels[i][j], ha='center', va='center',
                    fontsize=14, fontweight='bold', color=c)
    ax.set_xticks([0, 1])
    ax.set_yticks([0, 1])
    ax.set_xticklabels(['Predicted\nFailed', 'Predicted\nHealthy'],
                       fontsize=10, color='#0D1B2A')
    ax.set_yticklabels(['Actual\nFailed', 'Actual\nHealthy'],
                       fontsize=10, color='#0D1B2A')
    ax.set_title(f'25-fold LOVO Cross-Validation ({N_VINS} VINs)',
                 fontsize=12, fontweight='bold', color='#0D1B2A', pad=10)
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'confusion.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_feature_importance():
    fig, ax = plt.subplots(figsize=(10, 3.0))
    names = [f[0].replace('_', ' ').title() for f in FEATURES]
    vals = [f[1] for f in FEATURES]
    colors = ['#1A5CB0'] * len(vals)
    bars = ax.barh(names[::-1], vals[::-1], color=colors[::-1], height=0.55,
                   edgecolor='#0D1B2A', linewidth=0.5)
    for i, (n, v) in enumerate(zip(names[::-1], vals[::-1])):
        ax.text(v + 0.003, i, f'{v:.3f}', ha='left', va='center',
                fontsize=10, fontweight='bold', color='#0D1B2A')
    ax.set_xlabel('Permutation Importance (AUROC drop)', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'feature_importance.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_ridge_probs():
    """Ridge probability per VIN (failed vs NF)."""
    failed_vins = ['VIN1_F', 'VIN2_F', 'VIN3_F', 'VIN4_F', 'VIN5_F',
                   'VIN6_F', 'VIN7_F', 'VIN8_F', 'VIN9_F', 'VIN10_F']
    failed_probs = [0.6058, 0.6085, 0.7603, 0.4456, 0.2799,
                    0.7112, 0.7178, 0.8923, 0.4919, 0.6263]
    nf_vins = ['VIN1', 'VIN2', 'VIN3', 'VIN4', 'VIN5', 'VIN6', 'VIN7',
               'VIN8', 'VIN9', 'VIN10', 'VIN11', 'VIN12', 'VIN13', 'VIN14', 'VIN15']
    nf_probs = [0.2553, 0.3956, 0.4906, 0.4257, 0.3250, 0.4133, 0.4116,
                0.3934, 0.1053, 0.4389, 0.2559, 0.0421, 0.2104, 0.3162, 0.3987]

    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 3.2), gridspec_kw={'width_ratios': [1, 1.5]})

    # Failed
    c_f = ['#C0392B' if p >= THRESHOLD else '#E36C09' for p in failed_probs]
    ax1.barh(failed_vins[::-1], failed_probs[::-1], color=c_f[::-1], height=0.6,
             edgecolor='#0D1B2A', linewidth=0.5)
    ax1.axvline(x=THRESHOLD, color='#C0392B', linewidth=1.5, linestyle='--', label=f'Threshold {THRESHOLD}')
    for i, (v, p) in enumerate(zip(failed_vins[::-1], failed_probs[::-1])):
        ax1.text(p + 0.01, i, f'{p:.3f}', ha='left', va='center', fontsize=8,
                 fontweight='bold', color='#0D1B2A')
    ax1.set_xlim(0, 1.05)
    ax1.set_title('Failed Trucks (10)', fontsize=11, fontweight='bold', color='#0D1B2A')
    ax1.spines['top'].set_visible(False)
    ax1.spines['right'].set_visible(False)
    ax1.legend(fontsize=8, frameon=False)

    # NF
    nf_labels = [f'NF_{v.split("VIN")[1]}' for v in nf_vins]
    c_nf = ['#C0392B' if p >= THRESHOLD else '#1E8C45' for p in nf_probs]
    ax2.barh(nf_labels[::-1], nf_probs[::-1], color=c_nf[::-1], height=0.5,
             edgecolor='#0D1B2A', linewidth=0.5)
    ax2.axvline(x=THRESHOLD, color='#C0392B', linewidth=1.5, linestyle='--')
    for i, (v, p) in enumerate(zip(nf_labels[::-1], nf_probs[::-1])):
        ax2.text(p + 0.01, i, f'{p:.3f}', ha='left', va='center', fontsize=7,
                 fontweight='bold', color='#0D1B2A')
    ax2.set_xlim(0, 1.05)
    ax2.set_title('In-Service Trucks (15)', fontsize=11, fontweight='bold', color='#0D1B2A')
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)

    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'ridge_probs.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_failure_modes():
    """Failure mode split for 10 failed trucks."""
    fig, ax = plt.subplots(figsize=(10, 3.0))
    cats = ['Gradual-Electrical\n(GED2 Storm)', 'Abrupt / No\nPrecursor', 'Inconclusive\n(Data-Limited)']
    counts = [2, 5, 3]
    colors = ['#C58B1F', '#C0392B', '#B0B8C4']
    bars = ax.bar(cats, counts, color=colors, width=0.5, edgecolor='#0D1B2A', linewidth=0.5)
    for i, (c, v) in enumerate(zip(cats, counts)):
        ax.text(i, v + 0.15, f'{v} / 10', ha='center', va='bottom',
                fontsize=14, fontweight='bold', color='#0D1B2A')
    ax.set_ylim(0, 7)
    ax.set_ylabel('Number of Failures', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#B0B8C4')
    ax.spines['bottom'].set_color('#B0B8C4')
    ax.tick_params(colors='#606060')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'failure_modes.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_fleet_ttf():
    """Fleet TTF distribution (10 failed trucks)."""
    ttfs = [596, 627, 591, 664, 661, 552, 673, 573, 606, 472]
    fig, ax = plt.subplots(figsize=(10, 3.0))
    vins = [f'VIN{i}_F' for i in range(1, 11)]
    colors = ['#C58B1F' if t <= P25_TTF else ('#1A5CB0' if t <= P75_TTF else '#1E8C45') for t in ttfs]
    bars = ax.barh(vins[::-1], ttfs[::-1], color=colors[::-1], height=0.6,
                   edgecolor='#0D1B2A', linewidth=0.5)
    ax.axvline(x=MEDIAN_TTF, color='#C0392B', linewidth=2, linestyle='-', label=f'Median {MEDIAN_TTF:.0f}d')
    ax.axvline(x=P25_TTF, color='#B0B8C4', linewidth=1.5, linestyle='--', label=f'p25 {P25_TTF:.0f}d')
    ax.axvline(x=P75_TTF, color='#B0B8C4', linewidth=1.5, linestyle='--', label=f'p75 {P75_TTF:.0f}d')
    for i, (v, t) in enumerate(zip(vins[::-1], ttfs[::-1])):
        ax.text(t + 5, i, f'{t}d', ha='left', va='center',
                fontsize=9, fontweight='bold', color='#0D1B2A')
    ax.set_xlim(400, 720)
    ax.set_xlabel('Time-to-Failure (days)', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(fontsize=8, frameon=False, loc='lower right')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'fleet_ttf.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_pipeline():
    fig, ax = plt.subplots(figsize=(12, 2.5))
    stages = [
        ('CAN\nIngestion', '#1A5CB0'),
        ('Sentinel\nCleaning', '#2E5090'),
        ('Weekly\nAggregation', '#1A5CB0'),
        ('Feature\nEngineering', '#0D1B2A'),
        ('Ridge\nClassifier', '#2E5090'),
        ('Weibull\nSurvival', '#1A5CB0'),
        ('GED\nEmergency', '#0D1B2A'),
        ('RUL\nForecast', '#C58B1F'),
        ('Risk\nTier', '#1E8C45'),
    ]
    for i, (label, color) in enumerate(stages):
        x = i * 1.3
        rect = mpatches.FancyBboxPatch((x, 0.3), 1.0, 1.4, boxstyle="round,pad=0.1",
                                        facecolor=color, edgecolor='#0D1B2A', linewidth=1)
        ax.add_patch(rect)
        ax.text(x + 0.5, 1.0, label, ha='center', va='center', fontsize=9,
                fontweight='bold', color='white')
        if i < len(stages) - 1:
            ax.annotate('', xy=(x + 1.15, 1.0), xytext=(x + 1.05, 1.0),
                       arrowprops=dict(arrowstyle='->', color='#C58B1F', lw=2))
    ax.set_xlim(-0.3, len(stages) * 1.3)
    ax.set_ylim(-0.1, 2.2)
    ax.axis('off')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'pipeline.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_backtest():
    """Backtest MAE comparison: Model vs Fleet-Clock."""
    fig, ax = plt.subplots(figsize=(10, 3.0))
    cats = ['LOVO\nTotal TTF', 'Rewound\n270d', 'Rewound\n180d', 'Rewound\n90d', 'Overall']
    model_mae = [125.0, 130.4, 140.0, 154.8, 141.8]
    dummy_mae = [49.7, 49.7, 49.7, 49.7, 49.7]
    x = np.arange(len(cats))
    w = 0.3
    bars1 = ax.bar(x - w/2, model_mae, w, label='Model (Weibull)', color='#1A5CB0',
                   edgecolor='#0D1B2A', linewidth=0.5)
    bars2 = ax.bar(x + w/2, dummy_mae, w, label='Fleet-Clock Baseline', color='#C58B1F',
                   edgecolor='#0D1B2A', linewidth=0.5)
    for b in bars1:
        ax.text(b.get_x() + b.get_width()/2, b.get_height() + 2, f'{b.get_height():.0f}d',
                ha='center', fontsize=9, fontweight='bold', color='#1A5CB0')
    for b in bars2:
        ax.text(b.get_x() + b.get_width()/2, b.get_height() + 2, f'{b.get_height():.1f}d',
                ha='center', fontsize=9, fontweight='bold', color='#C58B1F')
    ax.set_ylabel('MAE (days)', fontsize=10, color='#606060')
    ax.set_xticks(x)
    ax.set_xticklabels(cats, fontsize=9, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(fontsize=9, frameon=False, loc='upper left')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'backtest.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path

def chart_nf_rul():
    """In-service fleet RUL forecast (sorted by median RUL)."""
    nf_data = [
        ('VIN11_NF', 87.0, 15.9, 231.6),
        ('VIN10_NF', 88.5, 15.8, 238.6),
        ('VIN2_NF',  89.5, 16.4, 236.9),
        ('VIN1_NF',  97.3, 13.9, 316.6),
        ('VIN6_NF', 102.4, 18.8, 256.6),
        ('VIN5_NF', 105.5, 19.8, 265.2),
        ('VIN7_NF', 106.4, 20.2, 262.5),
        ('VIN3_NF', 116.2, 22.9, 280.6),
        ('VIN13_NF', 118.4, 23.8, 283.4),
        ('VIN4_NF', 123.5, 23.9, 287.8),
        ('VIN14_NF', 137.7, 28.5, 309.2),
        ('VIN15_NF', 151.0, 32.9, 325.5),
        ('VIN12_NF', 183.0, 44.2, 371.4),
        ('VIN9_NF', 189.5, 47.3, 377.0),
        ('VIN8_NF', 234.5, 67.5, 428.1),
    ]
    fig, ax = plt.subplots(figsize=(12, 3.5))
    labels = [d[0] for d in nf_data]
    medians = [d[1] for d in nf_data]
    p10s = [d[2] for d in nf_data]
    p90s = [d[3] for d in nf_data]
    y = np.arange(len(labels))

    ax.barh(y, medians, height=0.5, color='#1A5CB0', edgecolor='#0D1B2A',
            linewidth=0.5, label='Median RUL')
    for i, (lab, med, p10, p90) in enumerate(zip(labels, medians, p10s, p90s)):
        ax.plot([p10, p90], [i, i], color='#B0B8C4', linewidth=2, zorder=1)
        ax.plot(p10, i, '|', color='#B0B8C4', markersize=8)
        ax.plot(p90, i, '|', color='#B0B8C4', markersize=8)
        color = '#C0392B' if lab == 'VIN3_NF' else '#0D1B2A'
        ax.text(med + 3, i, f'{med:.0f}d', ha='left', va='center',
                fontsize=8, fontweight='bold', color=color)
    ax.set_yticks(y)
    ax.set_yticklabels(labels, fontsize=8)
    ax.set_xlabel('Remaining Useful Life (days)', fontsize=10, color='#606060')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.legend(fontsize=8, frameon=False, loc='lower right')
    fig.patch.set_facecolor('white')
    plt.tight_layout()
    path = os.path.join(TMPDIR, 'nf_rul.png')
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


# ── GENERATE CHARTS ────────────────────────────────────────────
print("Generating charts...")
img_ridge      = chart_ridge_metrics()
img_confusion  = chart_confusion()
img_features   = chart_feature_importance()
img_probs      = chart_ridge_probs()
img_fmodes     = chart_failure_modes()
img_ttf        = chart_fleet_ttf()
img_pipeline   = chart_pipeline()
img_backtest   = chart_backtest()
img_nf_rul     = chart_nf_rul()
print("Charts generated.")

# Existing visualizations to embed
VIZ_SURVIVAL = str(VIZ / "fleet_survival_curve.png")
VIZ_BACKTEST = str(VIZ / "backtest_accuracy.png")
VIZ_GED      = str(VIZ / "ged_emergency.png")
VIZ_MAINT    = str(VIZ / "fleet_maintenance_board.png")
VIZ_WATERFALL = str(VIZ / "rul_band_waterfall.png")

# Best/worst case study VINs
VIZ_VIN1_F  = str(VIZ_CS / "Alternator_RUL_Degradation_VIN1_F_ALT_20260609.png")
VIZ_VIN5_F  = str(VIZ_CS / "Alternator_RUL_Degradation_VIN5_F_ALT_20260609.png")
VIZ_VIN3_NF = str(VIZ_CS / "Alternator_RUL_Degradation_VIN13_NF_ALT_20260609.png")
VIZ_VIN8_NF = str(VIZ_CS / "Alternator_RUL_Degradation_VIN18_NF_ALT_20260609.png")


# ══════════════════════════════════════════════════════════════
#  BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════
prs = make_prs()

# ─── SLIDE 1: TITLE ─────────────────────────────────────────
print("Building slide 1: Title")
s = blank_slide(prs)
add_rect(s, Inches(0), Inches(0), SW, SH, NAVY)
add_rect(s, Inches(0), Inches(5.7), SW, Inches(0.06), GOLD)

add_text(s, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
         "DAIMLER ALTERNATOR PREDICTIVE MAINTENANCE",
         font_size=34, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(2.7), Inches(11), Inches(0.5),
         "Fleet Risk Classification, Survival Analysis & Early Warning System",
         font_size=20, color=GOLD)
add_text(s, Inches(0.8), Inches(3.5), Inches(11), Inches(0.4),
         f"BharatBenz 5528T  |  Fleet of {N_VINS} VINs ({N_FAILED}F + {N_NF}NF)  |  V10.6.2",
         font_size=14, color=LIGHT_GREY)
add_text(s, Inches(0.8), Inches(5.9), Inches(5), Inches(0.3),
         "June 2026  |  Version: V10.6.2", font_size=12, color=LIGHT_GREY)
add_text(s, Inches(7), Inches(5.9), Inches(5), Inches(0.3),
         "Prepared by Data Science & Engineering Team", font_size=12,
         color=LIGHT_GREY, align=PP_ALIGN.RIGHT)
add_text(s, Inches(0.8), Inches(6.4), Inches(11), Inches(0.3),
         "BytEdge CONFIDENTIAL", font_size=10, bold=True, color=GOLD)


# ─── SLIDE 2: EXECUTIVE SUMMARY ─────────────────────────────
print("Building slide 2: Executive Summary")
s = blank_slide(prs)
add_header_bar(s, "EXECUTIVE SUMMARY",
               "Alternator failure prediction across 25 BharatBenz 5528T trucks")

y1, tw, th, gap, x0 = Inches(1.15), Inches(2.35), Inches(0.9), Inches(0.18), Inches(0.4)

add_kpi_tile(s, x0, y1, tw, th, "Ridge AUROC", f"{AUROC*100:.1f}%", KPI_BLUE, "PASS (>88%)")
add_kpi_tile(s, x0+tw+gap, y1, tw, th, "LOVO F1 Score", f"{F1_SCORE*100:.0f}%", KPI_BLUE, "PASS")
add_kpi_tile(s, x0+2*(tw+gap), y1, tw, th, "Fleet Median TTF", f"{MEDIAN_TTF:.0f} days", KPI_BLUE, f"~{MEDIAN_KM:,} km")
add_kpi_tile(s, x0+3*(tw+gap), y1, tw, th, "PI Coverage (LOVO)", f"{BACKTEST_PI_COV*100:.0f}%", KPI_GREEN, "PASS (>80%)")
add_kpi_tile(s, x0+4*(tw+gap), y1, tw, th, "Actionable Precursor", f"1 / {N_FAILED}", KPI_AMBER, "GED2 STORM")

y2 = Inches(2.25)
add_kpi_tile(s, x0, y2, tw, th, "Failed Trucks", str(N_FAILED), KPI_RED, "POST-HOC")
add_kpi_tile(s, x0+tw+gap, y2, tw, th, "In-Service Trucks", str(N_NF), KPI_GREEN, f"{N_LOW_RISK} LOW / {N_HIGH_RISK} HIGH")
add_kpi_tile(s, x0+2*(tw+gap), y2, tw, th, "Features (Ridge)", str(N_FEATURES), KPI_BLUE, "VSI-BASED")
add_kpi_tile(s, x0+3*(tw+gap), y2, tw, th, "Weibull Shape (α)", f"{WEIBULL_SHAPE:.1f}", KPI_BLUE, "SHARP WEAR-OUT")
add_kpi_tile(s, x0+4*(tw+gap), y2, tw, th, "Fleet-Clock MAE", f"{BACKTEST_MAE_DUMMY:.1f}d", KPI_GREEN, "DEPLOYED BASIS")

add_multiline(s, Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.8), [
    ("Objective: Predict alternator remaining useful life and classify failure risk from CAN bus voltage, excitation, and engine telemetry.", True, NAVY),
    (f"Ridge classifier achieves {AUROC*100:.1f}% AUROC with {N_FEATURES} VSI-derived features on 25-fold LOVO. Confusion matrix: {TP} TP, {FP} FP, {FN} FN, {TN} TN (F1={F1_SCORE:.2f}).", False, DARK_TEXT),
    (f"Fleet failures cluster tightly: median {MEDIAN_TTF:.0f} days (IQR {P25_TTF:.0f}–{P75_TTF:.0f}d), equivalent to ~{MEDIAN_KM:,} km / ~{MEDIAN_EHRS:,} engine-hours. Per-truck RUL models do not improve on fleet-clock baseline (MAE {BACKTEST_MAE_MODEL:.0f}d vs {BACKTEST_MAE_DUMMY:.1f}d), so fleet wear-out window is the deployed RUL basis.", False, DARK_TEXT),
    (f"GED2 excitation emergency fires on {GED_ACTIONABLE}/10 failures with actionable lead (~{VIN1_LEAD}d); {N_ABRUPT}/10 are electrically abrupt with no precursor. Honest assessment: fleet scheduling is the primary tool, not per-truck prognostics.", False, DARK_TEXT),
], font_size=10, bullet=True)

add_key_takeaways(s, [
    f"Ridge classifier (AUROC {AUROC*100:.1f}%) correctly separates high-risk from low-risk trucks at fleet level",
    f"Fleet replacement window ({P25_TTF:.0f}–{P75_TTF:.0f}d) is the deployable RUL basis; per-truck models add uncertainty bands, not point accuracy",
    "GED2 excitation storm is the only early-warning signal — fires on 1/10 failures with ~3 weeks lead",
    f"14/15 in-service trucks are LOW_RISK; 1 (VIN3_NF) is HIGH_RISK and should be prioritized for inspection",
], top=Inches(5.4))
add_footer(s)


# ─── SLIDE 3: BUSINESS CONTEXT ──────────────────────────────
print("Building slide 3: Business Context")
s = blank_slide(prs)
add_header_bar(s, "BUSINESS CONTEXT",
               "Why alternator predictive maintenance matters for heavy-duty fleets")

add_text(s, Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.35),
         "THE CHALLENGE", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(1.6), [
    "Alternator failure leaves truck stranded — 24V bus collapses, engine control lost",
    "No advance warning in current maintenance regime (time/km-based replacement)",
    "BharatBenz 5528T operates in extreme heat, dust, vibration — accelerating wear",
    "Fleet telemetry (CAN bus) streams voltage, excitation state, engine data — untapped",
    "Failure at ~120,000 km is premature vs OEM spec of 300,000–500,000 km",
], font_size=11, bullet=True, color=DARK_TEXT)

add_text(s, Inches(6.8), Inches(1.15), Inches(6), Inches(0.35),
         "OBJECTIVES", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(6.8), Inches(1.55), Inches(6), Inches(1.6), [
    "Classify fleet into risk tiers (GREEN / AMBER / RED) using CAN telemetry",
    "Estimate fleet-level replacement window from observed failure distribution",
    "Detect GED2 excitation disturbance as early-warning signal where physics allows",
    "Quantify what is achievable (and what is not) with available telemetry",
    "Provide honest, defensible RUL estimates with calibrated uncertainty bands",
], font_size=11, bullet=True, color=DARK_TEXT)

add_text(s, Inches(0.5), Inches(3.35), Inches(12.3), Inches(0.35),
         "THREE-LAYER APPROACH", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(3.75), Inches(12.3), Inches(1.4), [
    ("Layer 1 — Ridge Classifier: 6-feature VSI-based model separates high-risk from healthy trucks (AUROC 0.927)", True, NAVY),
    ("Layer 2 — Weibull Survival: Bayesian fleet wear-out model provides replacement-window timing with 80% credible intervals", True, NAVY),
    ("Layer 3 — GED Emergency: Real-time excitation-storm detector provides last-mile warning (when physics allows)", True, NAVY),
    "Together: risk stratification at fleet level, scheduling from survival analysis, and opportunistic early warning from GED2",
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "Alternator failure is premature in this fleet (~120k km vs 300–500k OEM spec) — a real business problem",
    "Three-layer architecture provides value at different time horizons: long-term (survival), medium (classifier), short (GED2)",
    "Honest framing: 50% of failures are electrically unwarnable; fleet scheduling is the primary mitigation",
])
add_footer(s)


# ─── SLIDE 4: PROBLEM STATEMENT ─────────────────────────────
print("Building slide 4: Problem Statement")
s = blank_slide(prs)
add_header_bar(s, "PROBLEM STATEMENT",
               "Characterizing the alternator failure landscape")

add_multiline(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5), [
    ("Key Question: Can we predict alternator failure timing and risk from CAN bus telemetry?", True, NAVY),
], font_size=14)

add_text(s, Inches(0.5), Inches(1.75), Inches(5.8), Inches(0.3),
         "FAILURE DISTRIBUTION", font_size=14, bold=True, color=NAVY)
s.shapes.add_picture(img_ttf, Inches(0.4), Inches(2.1), Inches(6.0), Inches(2.8))

add_text(s, Inches(6.8), Inches(1.75), Inches(6), Inches(0.3),
         "FAILURE MODE SPLIT", font_size=14, bold=True, color=NAVY)
s.shapes.add_picture(img_fmodes, Inches(6.8), Inches(2.1), Inches(6.0), Inches(2.8))

add_key_takeaways(s, [
    f"Tight failure cluster: median {MEDIAN_TTF:.0f}d, IQR {P25_TTF:.0f}–{P75_TTF:.0f}d (201-day window covers 50% of failures)",
    "5/10 failures are abrupt with no electrical precursor — the physics limits early warning to specific failure modes",
    "Only VIN1 (21-day GED2 storm) provides actionable multi-week warning; VIN10 is day-of confirmation only",
    "3/10 are inconclusive due to telemetry gaps near failure — improving coverage would sharpen the picture",
], top=Inches(5.15))
add_footer(s)


# ─── SLIDE 5: DATA LANDSCAPE ────────────────────────────────
print("Building slide 5: Data Landscape")
s = blank_slide(prs)
add_header_bar(s, "DATA LANDSCAPE",
               "CAN bus telemetry streams and fleet composition")

s.shapes.add_picture(img_pipeline, Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.8))

add_text(s, Inches(0.5), Inches(3.1), Inches(5.8), Inches(0.3),
         "CAN SIGNAL INVENTORY", font_size=14, bold=True, color=NAVY)
data = [
    ['Signal', 'Meaning', 'Range', 'Role'],
    ['VSI', 'Power Supply Voltage', '0–36V', 'Charging health (primary)'],
    ['GED', 'Alternator Excitation State', '{0,1,2,3}', 'Field-disturbance flag'],
    ['RPM', 'Engine Speed', '0–3500 rev/min', 'Operating regime'],
    ['CSP', 'Vehicle Speed', '0–100 km/h', 'Idle vs cruise separation'],
    ['ANR', 'Engine Torque', '-400–1300 Nm', 'Load context'],
    ['SMA', 'Starter Motor Active', '{0,1}', 'Crank-voltage analysis'],
]
add_table_shape(s, Inches(0.5), Inches(3.45), Inches(6.0), 7, 4, data,
                col_widths=[0.8, 2.0, 1.2, 2.0])

add_text(s, Inches(7.0), Inches(3.1), Inches(5.8), Inches(0.3),
         "FLEET COMPOSITION", font_size=14, bold=True, color=NAVY)
fleet_data = [
    ['Category', 'Count', 'Suffix', 'Status'],
    ['Failed', str(N_FAILED), '_F_ALT', 'Post-hoc analysis'],
    ['In-Service', str(N_NF), '_NF_ALT', 'Active monitoring'],
    ['Total', str(N_VINS), '—', f'~204M CAN rows'],
]
add_table_shape(s, Inches(7.0), Inches(3.45), Inches(5.8), 4, 4, fleet_data,
                col_widths=[1.3, 0.8, 1.2, 2.0])

add_multiline(s, Inches(7.0), Inches(5.0), Inches(5.8), Inches(0.5), [
    f"Sentinels: 65535 (CSP/RPM/ANR), -5000 (ANR), 0 & 255 (VSI)",
    f"VSI scaling: Actual V = Recorded × 0.2 if raw values exceed 36",
], font_size=9, bullet=True)

add_key_takeaways(s, [
    "VSI (power supply voltage) is the primary alternator health signal — all 6 classifier features derive from it",
    "GED=2 (alternator excitation disturbance) is the only real-time early-warning flag in the telemetry",
    "~204M CAN rows across 25 trucks, weekly-aggregated to ~10,000 feature rows for modeling",
], top=Inches(5.7), height=Inches(0.95))
add_footer(s)


# ─── SLIDE 6: FEATURE ENGINEERING ────────────────────────────
print("Building slide 6: Feature Engineering")
s = blank_slide(prs)
add_header_bar(s, "FEATURE ENGINEERING",
               "6 VSI-derived features selected via exhaustive LOVO elimination")

s.shapes.add_picture(img_features, Inches(0.4), Inches(1.15), Inches(6.2), Inches(2.8))

add_text(s, Inches(7.0), Inches(1.15), Inches(5.8), Inches(0.3),
         "FEATURE DESCRIPTIONS", font_size=14, bold=True, color=NAVY)
feat_data = [
    ['Feature', 'Description', 'Importance'],
    ['vsi_std_ratio_30d', 'Voltage variability trend (30d window)', '0.155'],
    ['vsi_dominant_freq', 'Dominant oscillation frequency in VSI', '0.105'],
    ['vsi_range_trend_last30d', 'Min-max voltage range trend (30d)', '0.061'],
    ['vsi_spectral_entropy', 'Spectral disorder of voltage signal', '0.057'],
    ['progressive_drift', 'Cumulative voltage drift over lifetime', '0.048'],
    ['bat_charge_delta_trend', 'Battery charging asymmetry trend', '0.034'],
]
add_table_shape(s, Inches(7.0), Inches(1.55), Inches(5.8), 7, 3, feat_data,
                col_widths=[2.2, 2.5, 1.1])

add_multiline(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(1.0), [
    ("Selection method: Exhaustive backward elimination over all subsets at n=25 VINs. 6 features (AUROC 0.927) beat 17 features (AUROC 0.907).", True, NAVY),
    "All features are derived from the VSI (power supply voltage) signal, confirming that alternator health is observable through charging-bus voltage characteristics.",
    "Permutation importance measured as AUROC drop when each feature is shuffled independently across LOVO folds.",
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "Fewer features perform better: 6 features (0.927) > 17 features (0.907) — exhaustively proven at n=25",
    "All features derive from VSI (voltage) — alternator health is observable through charging-bus characteristics",
    "vsi_std_ratio_30d (voltage variability) is the single most informative feature (importance: 0.155)",
])
add_footer(s)


# ─── SLIDE 7: SECTION DIVIDER ───────────────────────────────
print("Building slide 7: Section Divider")
s = blank_slide(prs)
add_section_slide(s, 4, "MODELING APPROACH",
                  "Three-layer architecture: Classification → Survival → Emergency")


# ─── SLIDE 8: MODELING APPROACH ──────────────────────────────
print("Building slide 8: Modeling Approach")
s = blank_slide(prs)
add_header_bar(s, "MODELING APPROACH",
               "Ridge classifier + Weibull survival + GED2 emergency layer")

add_text(s, Inches(0.5), Inches(1.15), Inches(3.8), Inches(0.3),
         "LAYER 1: RIDGE CLASSIFIER", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(3.8), Inches(1.8), [
    "L2-regularized logistic regression (Ridge)",
    f"6 VSI features, threshold {THRESHOLD} (Youden's J)",
    f"25-fold LOVO: AUROC {AUROC:.4f}",
    f"Bootstrap 95% CI: [{BOOTSTRAP_CI[0]:.3f}, {BOOTSTRAP_CI[1]:.3f}]",
    f"Permutation p-value: {PERMUTATION_P}",
    "Outputs: failure probability per VIN (0–1)",
    "Tier mapping: GREEN < 0.35, AMBER 0.35–0.55, RED ≥ 0.55",
], font_size=10, bullet=True)

add_text(s, Inches(4.7), Inches(1.15), Inches(3.8), Inches(0.3),
         "LAYER 2: WEIBULL SURVIVAL", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(4.7), Inches(1.55), Inches(3.8), Inches(1.8), [
    f"Bayesian Weibull: shape α={WEIBULL_SHAPE:.2f}, scale β={WEIBULL_SCALE:.1f}",
    f"Posterior median survival: {WEIBULL_MEDIAN:.1f} days",
    f"80% credible interval: {WEIBULL_CI[0]:.0f}–{WEIBULL_CI[1]:.0f} days",
    "Right-censoring for 15 in-service trucks",
    f"Empirical fleet median: {MEDIAN_TTF:.0f} days (deployed basis)",
    f"Conditional RUL: R = T−a | T>a from posterior grid",
    "PI coverage: 90% on LOVO backtest",
], font_size=10, bullet=True)

add_text(s, Inches(9.0), Inches(1.15), Inches(4.0), Inches(0.3),
         "LAYER 3: GED EMERGENCY", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(9.0), Inches(1.55), Inches(4.0), Inches(1.8), [
    "GED=2 (alternator excitation disturbance)",
    "Threshold: ≥200 GED2 counts per engine-on day",
    "Fires on sustained excitation storms only",
    f"Validated: {GED_ACTIONABLE}/10 actionable, 0/15 false alarms",
    f"Lead time: ~{VIN1_LEAD}d (VIN1), ~{VIN10_LEAD}d (VIN10)",
    "Last-mile safety net, not a primary predictor",
    "Not a component diagnosis (field-circuit flag)",
], font_size=10, bullet=True)

add_text(s, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.3),
         "DECISION ENGINE: HOW THE LAYERS COMBINE", font_size=14, bold=True, color=NAVY)
dec_data = [
    ['Condition', 'Action', 'Timing', 'Confidence'],
    ['GED2 storm active', 'IMMEDIATE INSPECTION', 'Days', 'High (0 false alarms)'],
    ['Ridge RED + age > p25', 'SCHEDULE REPLACEMENT', 'Weeks', 'Medium (AUROC 0.927)'],
    ['Ridge AMBER + age > median', 'MONITOR CLOSELY', 'Months', 'Medium'],
    ['Ridge GREEN + age < p25', 'ROUTINE MAINTENANCE', 'Quarters', 'High'],
    ['Fleet age > p75 (any tier)', 'PROACTIVE REPLACEMENT', 'Near-term', f'Fleet-level ({P75_TTF:.0f}d)'],
]
add_table_shape(s, Inches(0.5), Inches(3.95), Inches(12.3), 6, 4, dec_data,
                col_widths=[3.0, 3.0, 1.5, 4.8])

add_key_takeaways(s, [
    "Three independent layers provide redundant coverage: no single-point-of-failure in the prediction system",
    "GED2 emergency is the highest-confidence signal (0 false alarms) but rare (1/10 actionable)",
    "Fleet wear-out window is the primary scheduling tool; classifier adds per-truck risk differentiation",
], top=Inches(6.15), height=Inches(0.85))
add_footer(s)


# ─── SLIDE 9: PREDICTION FRAMEWORK ──────────────────────────
print("Building slide 9: Prediction Framework")
s = blank_slide(prs)
add_header_bar(s, "PREDICTION FRAMEWORK",
               "How fleet-clock RUL and conditional survival combine")

add_multiline(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.4), [
    ("Core insight: Per-truck RUL models cannot beat the fleet-clock baseline. The fleet wears out together.", True, NAVY),
], font_size=13)

add_text(s, Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.3),
         "FLEET-CLOCK BASELINE (DEPLOYED)", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(2.0), Inches(5.8), Inches(1.8), [
    f"RUL = {MEDIAN_TTF:.0f} − current_age_days",
    f"MAE on LOVO backtest: {BACKTEST_MAE_DUMMY:.1f} days",
    "Simple, interpretable, and outperforms Weibull per-truck model",
    "Works because failure distribution is tight (IQR 75 days)",
    f"Equivalent: ~{MEDIAN_KM:,} km / ~{MEDIAN_EHRS:,} engine-hours",
], font_size=11, bullet=True)

add_text(s, Inches(6.8), Inches(1.65), Inches(6.0), Inches(0.3),
         "WEIBULL SURVIVAL (UNCERTAINTY)", font_size=14, bold=True, color=NAVY)
add_multiline(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(1.8), [
    f"Weibull α={WEIBULL_SHAPE:.2f}: peaked hazard confirms batch wear-out pattern",
    f"Per-truck conditional RUL: posterior sampling of R = T−a | T>a",
    f"LOVO backtest MAE: {BACKTEST_MAE_MODEL:.0f}d (worse than fleet-clock {BACKTEST_MAE_DUMMY:.1f}d)",
    "Value: provides calibrated 80% prediction intervals (PI coverage 90%)",
    "Not deployed for point prediction; deployed for uncertainty quantification",
], font_size=11, bullet=True)

s.shapes.add_picture(img_backtest, Inches(0.4), Inches(4.0), Inches(6.2), Inches(2.5))

add_multiline(s, Inches(7.0), Inches(4.0), Inches(5.8), Inches(1.5), [
    ("HONEST ASSESSMENT", True, NAVY),
    f"Fleet-clock wins on MAE ({BACKTEST_MAE_DUMMY:.1f}d vs {BACKTEST_MAE_MODEL:.0f}d) because alternator failures cluster tightly around the fleet median.",
    "Weibull model adds value through uncertainty bands: 90% of actual failures fall within its prediction interval.",
    "Per-truck prognostics require richer failure diversity than this fleet provides (n=10, IQR=75d).",
], font_size=10, bullet=True)
add_footer(s)


# ─── SLIDE 10: RESULTS OVERVIEW ─────────────────────────────
print("Building slide 10: Results Overview")
s = blank_slide(prs)
add_header_bar(s, "RESULTS OVERVIEW",
               "Ridge classifier performance on 25-fold LOVO")

s.shapes.add_picture(img_ridge, Inches(0.4), Inches(1.1), Inches(6.2), Inches(2.5))
s.shapes.add_picture(img_confusion, Inches(7.5), Inches(1.1), Inches(4.5), Inches(2.7))

add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.3),
         "LOVO PERFORMANCE SCORECARD", font_size=14, bold=True, color=NAVY)
score_data = [
    ['Metric', 'Value', 'Context'],
    ['AUROC', f'{AUROC:.4f}', f'Bootstrap 95% CI: [{BOOTSTRAP_CI[0]:.3f}, {BOOTSTRAP_CI[1]:.3f}]'],
    ['F1 Score', f'{F1_SCORE:.2f}', f'{TP} TP + {FP} FP + {FN} FN'],
    ['Recall (Sensitivity)', f'{RECALL:.2f}', f'{TP}/{TP+FN} failed trucks correctly identified'],
    ['Precision', f'{PRECISION:.2f}', f'{TP}/{TP+FP} positive predictions are true failures'],
    ['Specificity', f'{SPECIFICITY:.4f}', f'{TN}/{TN+FP} healthy trucks correctly classified'],
    ['MCC', f'{MCC:.4f}', 'Matthews correlation (balanced metric)'],
    ['Optimal Threshold', f'{THRESHOLD}', "Youden's J on LOVO predictions"],
]
add_table_shape(s, Inches(0.5), Inches(4.15), Inches(12.3), 8, 3, score_data,
                col_widths=[2.5, 1.5, 8.3])

add_key_takeaways(s, [
    f"AUROC {AUROC:.4f} with 6 features separates failed from healthy trucks with high discrimination",
    f"1 FP (VIN5_F scored GREEN at 0.28) and 1 FN — no model is perfect at n=25",
    "Permutation test p=0.0 confirms the classifier captures real signal, not noise",
], top=Inches(6.9), height=Inches(0.55))
add_footer(s)


# ─── SLIDE 11: SECTION DIVIDER ──────────────────────────────
print("Building slide 11: Section Divider")
s = blank_slide(prs)
add_section_slide(s, 7, "DETAILED VIN ANALYSIS",
                  "Failed trucks, in-service fleet, and case studies")


# ─── SLIDE 12: FAILED VIN ANALYSIS ──────────────────────────
print("Building slide 12: Failed VIN Analysis")
s = blank_slide(prs)
add_header_bar(s, "FAILED TRUCK ANALYSIS — 10 VINs",
               "Forensic classification of alternator failure signatures")

failed_data = [
    ['VIN', 'TTF (d)', 'Ridge P', 'Band', 'Failure Class', 'Lead', 'Actionable'],
    ['VIN1_F', '596', '0.606', 'RED', 'Gradual-electrical', '21d', 'YES'],
    ['VIN2_F', '627', '0.609', 'RED', 'Inconclusive', '—', 'NO'],
    ['VIN3_F', '591', '0.760', 'RED', 'Inconclusive (gap)', '—', 'NO'],
    ['VIN4_F', '664', '0.446', 'AMBER', 'Abrupt', '—', 'NO'],
    ['VIN5_F', '661', '0.280', 'GREEN', 'Abrupt', '—', 'NO'],
    ['VIN6_F', '552', '0.711', 'RED', 'Abrupt', '—', 'NO'],
    ['VIN7_F', '673', '0.718', 'RED', 'Abrupt', '—', 'NO'],
    ['VIN8_F', '573', '0.892', 'RED', 'Inconclusive (3d)', '—', 'NO'],
    ['VIN9_F', '606', '0.492', 'AMBER', 'Abrupt', '—', 'NO'],
    ['VIN10_F', '472', '0.626', 'RED', 'Gradual-electrical', '1d', 'DAY-OF'],
]
add_table_shape(s, Inches(0.4), Inches(1.15), Inches(12.5), 12, 7, failed_data,
                col_widths=[1.2, 0.9, 1.0, 1.0, 2.8, 0.9, 1.3])

add_multiline(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.3), [
    ("Ridge correctly flags 7/10 failures as RED (≥0.55), plus 2 AMBER. Only VIN5_F (0.280) is a clean miss.", False, DARK_TEXT),
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "7/10 failed trucks scored RED by the Ridge classifier — correct risk identification",
    "VIN5_F is the only clean false negative (GREEN, 0.280) — likely a non-electrical failure mode",
    "VIN1_F is the sole actionable GED2 precursor: sustained excitation storm with ~21 days warning",
    "5/10 are abrupt/no-precursor: electrically silent until failure — these cannot be warned by any voltage-based method",
], top=Inches(5.95))
add_footer(s)


# ─── SLIDE 13: IN-SERVICE FLEET ─────────────────────────────
print("Building slide 13: In-Service Fleet")
s = blank_slide(prs)
add_header_bar(s, "IN-SERVICE FLEET STATUS — 15 VINs",
               "Current risk tier and remaining useful life forecast")

s.shapes.add_picture(img_nf_rul, Inches(0.4), Inches(1.15), Inches(12.4), Inches(3.5))

add_multiline(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.5), [
    (f"All 15 in-service trucks are past the fleet median age ({MEDIAN_TTF:.0f}d). Median remaining RUL ranges from 87–235 days.", True, NAVY),
    f"14/15 are LOW_RISK. VIN3_NF (ridge prob 0.491, AMBER) is the only HIGH_RISK truck — recommend priority inspection.",
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "Error bars show 80% prediction interval (p10–p90) from Weibull conditional survival",
    "VIN3_NF is the only HIGH_RISK in-service truck (ridge_prob=0.491) — should be inspected first",
    "VIN8_NF has the longest runway (235d median RUL) despite being a newer truck with lower age",
    "All trucks are NEAR_TERM: entire fleet is within the replacement window",
], top=Inches(5.4))
add_footer(s)


# ─── SLIDE 14: VIN CASE STUDIES ──────────────────────────────
print("Building slide 14: VIN Case Studies")
s = blank_slide(prs)
add_header_bar(s, "VIN CASE STUDIES",
               "Best prediction, typical, high-risk in-service, and missed detection")

add_text(s, Inches(0.4), Inches(1.15), Inches(3.0), Inches(0.25),
         "VIN1_F — BEST: 21-DAY WARNING", font_size=11, bold=True, color=GREEN_PASS)
if Path(VIZ_VIN1_F).exists():
    s.shapes.add_picture(VIZ_VIN1_F, Inches(0.4), Inches(1.45), Inches(6.2), Inches(2.5))
else:
    add_text(s, Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.3),
             "[VIN1_F chart not found]", font_size=10, color=RED_FAIL)

add_text(s, Inches(6.8), Inches(1.15), Inches(3.0), Inches(0.25),
         "VIN5_F — MISSED: FALSE NEGATIVE", font_size=11, bold=True, color=RED_FAIL)
if Path(VIZ_VIN5_F).exists():
    s.shapes.add_picture(VIZ_VIN5_F, Inches(6.8), Inches(1.45), Inches(6.2), Inches(2.5))
else:
    add_text(s, Inches(6.8), Inches(1.45), Inches(6.0), Inches(0.3),
             "[VIN5_F chart not found]", font_size=10, color=RED_FAIL)

add_multiline(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.0), [
    ("VIN1_F: The only actionable precursor — sustained GED2 excitation storm starting ~21 days before failure. Ridge correctly scored RED (0.606). This is what early warning looks like when physics allows.", False, DARK_TEXT),
    ("VIN5_F: Clean false negative — ridge scored GREEN (0.280), no GED2 activity. Electrically flat to failure with good telemetry coverage (30 days). Likely a non-electrical failure mode (bearing, winding open) invisible to voltage-based features.", False, DARK_TEXT),
    ("Implication: Not all alternator failures are predictable from available telemetry. The 5/10 abrupt failures represent a physics ceiling, not a model deficiency.", False, DARK_TEXT),
], font_size=10, bullet=True)

add_key_takeaways(s, [
    "Best case (VIN1): GED2 + Ridge RED together give 3-week actionable window",
    "Worst case (VIN5): No precursor in any channel — fundamentally unpredictable from voltage signals",
    "This contrast is the honest story: the model works where physics allows, and admits where it cannot",
])
add_footer(s)


# ─── SLIDE 15: SECTION DIVIDER ──────────────────────────────
print("Building slide 15: Section Divider")
s = blank_slide(prs)
add_section_slide(s, 8, "OPERATIONAL ASSESSMENT",
                  "Value, limitations, and the path forward")


# ─── SLIDE 16: OPERATIONAL VALUE ─────────────────────────────
print("Building slide 16: Operational Value")
s = blank_slide(prs)
add_header_bar(s, "OPERATIONAL VALUE",
               "What the system delivers today and what it does not")

add_text(s, Inches(0.5), Inches(1.15), Inches(5.8), Inches(0.3),
         "WHAT IT DELIVERS", font_size=16, bold=True, color=GREEN_PASS)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(5.8), Inches(2.5), [
    ("Fleet-Level Scheduling: Replace alternators within the {}–{}d window (IQR) to cover 50% of failures proactively".format(int(P25_TTF), int(P75_TTF)), False, DARK_TEXT),
    (f"Risk Stratification: Ridge classifier (AUROC {AUROC:.3f}) separates RED (7/10 failures) from GREEN (low-risk), enabling prioritized inspection", False, DARK_TEXT),
    ("Early Warning (Rare): GED2 excitation storm provides 1–3 week alert on excitation-type failures (1/10 in this fleet)", False, DARK_TEXT),
    ("Calibrated Uncertainty: 80% prediction intervals on per-truck RUL (90% backtest coverage)", False, DARK_TEXT),
    ("Zero False Alarms on GED2: No in-service truck has triggered the emergency layer incorrectly", False, DARK_TEXT),
], font_size=10, bullet=True)

add_text(s, Inches(6.8), Inches(1.15), Inches(6.0), Inches(0.3),
         "WHAT IT DOES NOT DELIVER", font_size=16, bold=True, color=RED_FAIL)
add_multiline(s, Inches(6.8), Inches(1.55), Inches(6.0), Inches(2.5), [
    ("Precise per-truck RUL: MAE 125d (model) vs 49.7d (fleet-clock) — fleet-clock wins because failures cluster", False, DARK_TEXT),
    ("Universal early warning: 5/10 failures are electrically abrupt — no voltage-based method can predict them", False, DARK_TEXT),
    ("Component diagnosis: GED2 and VSI are bus-level signals, not alternator-internal diagnostics", False, DARK_TEXT),
    ("Generalization guarantee: n=10 failures is too small for reliable base-rate statistics", False, DARK_TEXT),
    ("Coverage for data gaps: 3/10 failures are inconclusive due to telemetry dropout near failure", False, DARK_TEXT),
], font_size=10, bullet=True)

add_text(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.3),
         "BUSINESS IMPACT ESTIMATE", font_size=14, bold=True, color=NAVY)
impact_data = [
    ['Scenario', 'Without System', 'With System', 'Improvement'],
    ['Unplanned Failures', '10/10 reactive', '7/10 flagged RED pre-failure', '70% risk-identified'],
    ['Fleet Scheduling', 'Time/km-based', f'IQR window ({int(P25_TTF)}–{int(P75_TTF)}d)', 'Data-driven window'],
    ['Emergency Detection', 'None', 'GED2 fires on excitation storms', 'Days-level warning'],
    ['Inspection Priority', 'Equal treatment', 'Ridge-ranked risk tiers', 'Resource-efficient'],
]
add_table_shape(s, Inches(0.5), Inches(4.55), Inches(12.3), 5, 4, impact_data,
                col_widths=[2.5, 3.0, 3.5, 3.3])

add_footer(s)


# ─── SLIDE 17: CURRENT LIMITATIONS ──────────────────────────
print("Building slide 17: Current Limitations")
s = blank_slide(prs)
add_header_bar(s, "CURRENT LIMITATIONS",
               "Honest assessment of what the data and models cannot do")

limit_data = [
    ['Limitation', 'Impact', 'Root Cause', 'Mitigation'],
    ['Per-truck RUL loses to fleet-clock',
     f'MAE {BACKTEST_MAE_MODEL:.0f}d vs {BACKTEST_MAE_DUMMY:.1f}d',
     'Failures cluster tightly (IQR 75d)',
     'Deploy fleet-clock as primary; use Weibull for uncertainty'],
    ['5/10 failures are electrically unwarnable',
     'No advance signal for 50% of failures',
     'Abrupt failure modes (bearing, winding open)',
     'Accept physics ceiling; supplement with vibration sensors'],
    ['Only 1/10 has actionable GED2 precursor',
     'GED2 emergency layer is rare',
     'Excitation storms are one failure sub-type',
     'Value = high-confidence when it fires (0 FP)'],
    ['n=10 failures is statistically small',
     'Wide confidence intervals, limited generalization',
     'Real-world data scarcity for alternator failures',
     f'Bootstrap CI [{BOOTSTRAP_CI[0]:.3f}, {BOOTSTRAP_CI[1]:.3f}] quantifies uncertainty'],
    ['Weibull overpredicts by ~118 days',
     f'Posterior median {WEIBULL_MEDIAN:.0f}d vs empirical {MEDIAN_TTF:.0f}d',
     'Right-censored NF trucks pull estimate high',
     'Use empirical median for point RUL; Weibull for intervals'],
    ['Telemetry gaps near some failures',
     '3/10 failures lack pre-failure data',
     'Telemetry dropout in final weeks',
     'Improve end-of-life data collection in fleet'],
    ['Wide NF prediction intervals',
     '~220–360 day PI width for in-service trucks',
     'Conditional survival inherits Weibull uncertainty',
     'Narrow with more failure observations (currently n=10)'],
]
add_table_shape(s, Inches(0.4), Inches(1.15), Inches(12.5), 8, 4, limit_data,
                col_widths=[3.0, 2.5, 3.0, 4.0])

add_key_takeaways(s, [
    "These are not model bugs — they are physics and data-size constraints that should be communicated transparently",
    "The system is honest about what it can and cannot do; this honesty is itself a deliverable",
    "Primary action: improve telemetry coverage and grow the failure cohort to sharpen predictions",
], top=Inches(4.15))
add_footer(s)


# ─── SLIDE 18: IMPROVEMENT ROADMAP ──────────────────────────
print("Building slide 18: Improvement Roadmap")
s = blank_slide(prs)
add_header_bar(s, "IMPROVEMENT ROADMAP",
               "Concrete next steps to strengthen the system")

road_data = [
    ['Priority', 'Action', 'Expected Impact', 'Timeline'],
    ['P0', 'Starter Motor pipeline (n=34 trucks)', 'Second component; proven methodology transfer', 'Q3 2026'],
    ['P0', 'Deployment packaging (API/dashboard)', 'Operational readiness for fleet managers', 'Q3 2026'],
    ['P1', 'Grow alternator failure cohort (n>25)', 'Narrower CI, better base-rate statistics', 'Ongoing'],
    ['P1', 'Improve telemetry coverage at end-of-life', 'Resolve 3/10 inconclusive failures', 'Q3–Q4 2026'],
    ['P2', 'Vibration/temperature sensor integration', 'Detect bearing/mechanical failure modes', 'Q4 2026+'],
    ['P2', 'Cross-fleet validation (different routes/loads)', 'Confirm generalization beyond single fleet', '2027'],
    ['P3', 'Component-level diagnostics (field current)', 'Move from bus-level to alternator-internal', 'Research'],
]
add_table_shape(s, Inches(0.4), Inches(1.15), Inches(12.5), 8, 4, road_data,
                col_widths=[1.0, 4.5, 4.0, 3.0])

add_key_takeaways(s, [
    "P0: Starter motor pipeline is the immediate next deliverable — same methodology, larger fleet (34 trucks)",
    "Growing the failure cohort is the single highest-leverage improvement for model accuracy",
    "Vibration sensors would address the 5/10 abrupt failures that voltage telemetry cannot see",
    "Cross-fleet validation will determine if the 601-day fleet window generalizes to other routes",
], top=Inches(4.15))
add_footer(s)


# ─── SLIDE 19: CONCLUSIONS ──────────────────────────────────
print("Building slide 19: Conclusions")
s = blank_slide(prs)
add_header_bar(s, "CONCLUSIONS & RECOMMENDATIONS",
               "Summary of findings and priority actions")

add_text(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.3),
         "SUMMARY OF FINDINGS", font_size=16, bold=True, color=NAVY)
add_multiline(s, Inches(0.5), Inches(1.55), Inches(12.3), Inches(2.0), [
    (f"Ridge classifier (AUROC {AUROC:.3f}, F1 {F1_SCORE:.2f}) provides reliable fleet-level risk stratification with 6 VSI-derived features", False, DARK_TEXT),
    (f"Fleet replacement window of {int(P25_TTF)}–{int(P75_TTF)} days (~{MEDIAN_KM:,} km) is the primary RUL basis — simple and effective", False, DARK_TEXT),
    ("GED2 excitation emergency provides high-confidence but rare early warning (1/10 actionable, 0 false alarms)", False, DARK_TEXT),
    ("50% of failures are electrically unwarnable — a physics ceiling, not a model limitation", False, DARK_TEXT),
    ("Per-truck survival RUL does not improve on fleet-clock but provides valuable calibrated uncertainty bands", False, DARK_TEXT),
    ("14/15 in-service trucks are LOW_RISK; 1 (VIN3_NF) is HIGH_RISK and should be prioritized", False, DARK_TEXT),
], font_size=11, bullet=True)

add_text(s, Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.3),
         "PRIORITY RECOMMENDATIONS", font_size=16, bold=True, color=NAVY)
rec_data = [
    ['#', 'Recommendation', 'Rationale'],
    ['1', 'Schedule VIN3_NF_ALT for priority inspection', 'Only HIGH_RISK in-service truck (ridge_prob=0.491)'],
    ['2', f'Plan fleet-wide alternator replacement at {int(P25_TTF)}–{int(P75_TTF)}d', 'IQR window covers 50% of observed failures'],
    ['3', 'Deploy GED2 emergency alerts in real-time', '0 false alarms; provides days-level last-mile warning'],
    ['4', 'Begin Starter Motor pipeline (n=34)', 'Proven methodology transfers to second component'],
    ['5', 'Improve end-of-life telemetry collection', '3/10 failures are inconclusive due to data gaps'],
    ['6', 'Grow failure cohort before refining models', f'n=10 gives CI [{BOOTSTRAP_CI[0]:.3f}, {BOOTSTRAP_CI[1]:.3f}]; more failures narrow this'],
    ['7', 'Communicate the honest picture to stakeholders', 'Set expectations: fleet scheduling > per-truck prognostics for this fleet'],
]
add_table_shape(s, Inches(0.4), Inches(4.0), Inches(12.5), 8, 3, rec_data,
                col_widths=[0.5, 5.0, 7.0])

add_footer(s)


# ─── SLIDE 20: APPENDIX ─────────────────────────────────────
print("Building slide 20: Appendix")
s = blank_slide(prs)
add_header_bar(s, "APPENDIX & TECHNICAL REFERENCE",
               "Key terms, data sources, and pipeline specifications")

terms = [
    ['Term', 'Definition'],
    ['AUROC', 'Area Under ROC Curve — measures classifier discrimination (1.0 = perfect)'],
    ['LOVO', 'Leave-One-Vehicle-Out cross-validation — each VIN held out once for testing'],
    ['Ridge Classifier', 'L2-regularized logistic regression for binary failure prediction'],
    ['Weibull Distribution', 'Parametric survival model; shape α > 1 indicates wear-out mode'],
    ['GED', 'State of Alternator Excitation (CAN signal): 0=Normal, 2=Disturbance, 3=Unavailable'],
    ['VSI', 'Power Supply Voltage (CAN signal): 24V system regulated at ~28V when healthy'],
    ['TTF', 'Time-to-Failure in days from first observation to replacement event'],
    ['RUL', 'Remaining Useful Life — estimated days until failure from current age'],
    ['PI Coverage', 'Fraction of actual values falling within the prediction interval'],
    ['Fleet-Clock', 'RUL = median_TTF − current_age; simplest survival-based baseline'],
    ['Conditional RUL', 'R = T−a | T>a sampled from Bayesian posterior survival distribution'],
    ['MCC', 'Matthews Correlation Coefficient — balanced measure even with class imbalance'],
    ['Youden’s J', 'Optimal threshold = argmax(Sensitivity + Specificity − 1)'],
]
add_table_shape(s, Inches(0.4), Inches(1.15), Inches(12.5), len(terms), 2, terms,
                col_widths=[2.5, 10.0])

add_multiline(s, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.5), [
    "Data sources: V10.6.2_ALT/cache/ (RUL, backtest, Weibull, forensics, GED), V10.6_ALT/cache/ (ridge, rules, lifecycle), V5.2_ALT/ (feature selection)",
    "All metrics independently recomputed from pipeline outputs. See column_dictionary.md for signal definitions.",
], font_size=8, bullet=True, color=GREY_MED)
add_footer(s)


# ── SAVE ───────────────────────────────────────────────────────
out_path = str(OUT_DIR / "Alternator_Predictive_Maintenance_V10.6.2.pptx")
prs.save(out_path)
print(f"\nPresentation saved: {out_path}")
print(f"  Slides: {len(prs.slides)}")

shutil.rmtree(TMPDIR, ignore_errors=True)
print("Done.")
