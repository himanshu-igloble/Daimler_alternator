"""V11_ALT_heuristics — technical presentation (python-pptx), V11-native."""
from __future__ import annotations

import importlib.util
import pathlib

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

_src = pathlib.Path(__file__).resolve().parents[1] / "src"


def _load(name):
    spec = importlib.util.spec_from_file_location(name, str(_src / f"{name}.py"))
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


cfg = _load("V11_ALT_heuristics_config")
FOR = cfg.FORENSICS
GRAPHS = cfg.V11_ROOT / "visualizations" / "V11_graphs"
NAVY, GOLD, GREEN = RGBColor(0x0D, 0x1B, 0x2A), RGBColor(0xC5, 0x8B, 0x1F), RGBColor(0x27, 0xAE, 0x60)


def _nums():
    comp = pd.read_csv(cfg.RESULTS_DIR / "V11_ALT_heuristics_comparison.csv")
    es = pd.read_csv(FOR / "earliest_signal_per_vin.csv")
    st = pd.read_csv(FOR / "nf_self_test.csv")
    cl = pd.read_csv(FOR / "compound_alarm_lovo.csv")
    return {
        "n11": int((es["verdict"] == "discriminative_precursor").sum()),
        "n1062": int((comp["v1062_horizon"].astype(str) != "none").sum()),
        "nf_false": int((st["verdict"] == "FALSE_ALARM").sum()),
        "cf": int(cl[(cl.group == "FAILED") & (cl.fired == True)].shape[0]),   # noqa: E712
        "cnf": int(cl[(cl.group == "NF") & (cl.fired == True)].shape[0]),      # noqa: E712
    }


def _title_slide(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12), Inches(2))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.size = Pt(20); p2.font.color.rgb = GOLD
    return s


def _bullets_slide(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    p = t.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = NAVY
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.6), Inches(5.4)).text_frame
    body.word_wrap = True
    for i, b in enumerate(bullets):
        par = body.paragraphs[0] if i == 0 else body.add_paragraph()
        par.text = "• " + b; par.font.size = Pt(18); par.font.color.rgb = NAVY
    return s


def _image_slide(prs, title, img):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    p = t.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(26); p.font.bold = True; p.font.color.rgb = NAVY
    if pathlib.Path(img).exists():
        s.shapes.add_picture(str(img), Inches(1.2), Inches(1.2), height=Inches(5.6))
    return s


def main():
    n = _nums()
    prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    _title_slide(prs, "Alternator Lead-Time Heuristics (V11)",
                 f"Recall {n['n11']}/10 vs V10.6.2 {n['n1062']}/10  |  {n['nf_false']}/15 false alarms")
    _bullets_slide(prs, "The problem V10.6.2 left open", [
        "Frozen classifier says WHICH truck (AUROC 0.927) but gives no timing.",
        "Per-truck RUL cannot beat a fleet-clock dummy (structural, n=10).",
        "Only genuine precursor was the GED=2 storm — fired for just 2/10 failures.",
        "Open question: can we get earlier, higher-recall precursors from the same 6 channels?"])
    _bullets_slide(prs, "Approach: 12 heuristics on 6 raw CAN channels", [
        "Load-normalised & regime-conditioned voltage (dVSI/dRPM, load residual, reg duty).",
        "Crank/recovery dynamics (post-crank recovery, crank effort).",
        "Excitation & instability (full GED states, idle hunting), sag typing, UV dose.",
        "Compound voting alarm (#11) + per-truck change-point (#12)."])
    _bullets_slide(prs, "Honest validation gate (unchanged from V10.6.2)", [
        "A deviation counts ONLY if within-truck z>=2 AND outside healthy-fleet p05-p95.",
        "MIN_EO_SAMPLES=200 trust filter; horizons 90/60/45/30/14/7 days.",
        "NF self-test: every healthy truck scored as-if-failing (false-alarm honesty).",
        "n=10 events — results are leads, not calibrated rates."])
    _image_slide(prs, "Result: lead-time recall", GRAPHS / "G1_recall_comparison.png")
    _image_slide(prs, "Per-truck head-to-head", GRAPHS / "G3_leadtime_dumbbell.png")
    _image_slide(prs, "Strongest new signal: post-crank recovery (#3) — episodic", GRAPHS / "G5_crank_recovery_trajectories.png")
    _image_slide(prs, "Which heuristics generalised", GRAPHS / "G2_feature_generalization.png")
    _image_slide(prs, "Compound early-watch alarm (#11)", GRAPHS / "G4_compound_alarm_leads.png")
    _image_slide(prs, "Change-point / resting voltage (exploratory)", GRAPHS / "G6_changepoint_resting.png")
    _bullets_slide(prs, "Limitations", [
        "Net gain is modest: +1 detection (VIN9), +1 earlier lead (VIN1).",
        "crank_recovery_t z-magnitudes inflated (healthy baseline ~0.05s); trust the 0/15 NF guard.",
        "The recovery signal is episodic (isolated spikes); VIN9's events span its whole life.",
        "VIN8 strict-gate hit rests on 3 trusted days (compound catches it at 90d).",
        "4/10 failures remain undetectable (abrupt/silent); no per-truck daily RUL implied."])
    _bullets_slide(prs, "Deployment recommendation", [
        "Add crank_recovery_t to the emergency channel alongside the GED=2 storm.",
        "Surface the compound 2-of-5 vote as an 'early-watch' tier (0 false alarms).",
        "Use sag-typing (high-load vs idle) + reg-duty for repair-direction guidance.",
        "WHICH (classifier) + WHEN-fleet (Weibull) unchanged from V10.6.2."])
    _bullets_slide(prs, "Appendix — data sources", [
        "All numbers from V11 cache/forensics + results/comparison.csv.",
        "No RUL/Weibull/backtest data (V11 is the precursor fork; see V10.6.2 for RUL).",
        "Graphs: visualizations/V11_graphs/.",
        "Pipeline reproducible via V11_ALT_heuristics_orchestrator.py."])
    out = cfg.V11_ROOT / "presentation" / "Alternator_LeadTime_Heuristics_V11.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"[v11 technical deck] wrote {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
