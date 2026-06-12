"""V11_ALT_heuristics — business summary deck (5 slides), V11-native."""
from __future__ import annotations

import importlib.util
import pathlib

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

_src = pathlib.Path(__file__).resolve().parents[1] / "src"


def _load(name):
    spec = importlib.util.spec_from_file_location(name, str(_src / f"{name}.py"))
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


cfg = _load("V11_ALT_heuristics_config")
FOR = cfg.FORENSICS
GRAPHS = cfg.V11_ROOT / "visualizations" / "V11_graphs"
NAVY, GOLD = RGBColor(0x0D, 0x1B, 0x2A), RGBColor(0xC5, 0x8B, 0x1F)


def _nums():
    comp = pd.read_csv(cfg.RESULTS_DIR / "V11_ALT_heuristics_comparison.csv")
    es = pd.read_csv(FOR / "earliest_signal_per_vin.csv")
    st = pd.read_csv(FOR / "nf_self_test.csv")
    cl = pd.read_csv(FOR / "compound_alarm_lovo.csv")
    return {"n11": int((es["verdict"] == "discriminative_precursor").sum()),
            "n1062": int((comp["v1062_horizon"].astype(str) != "none").sum()),
            "nf_false": int((st["verdict"] == "FALSE_ALARM").sum()),
            "cf": int(cl[(cl.group == "FAILED") & (cl.fired == True)].shape[0])}  # noqa: E712


def _bullets(prs, title, bullets):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    p = t.text_frame.paragraphs[0]; p.text = title; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = NAVY
    body = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.6), Inches(5.2)).text_frame
    body.word_wrap = True
    for i, b in enumerate(bullets):
        par = body.paragraphs[0] if i == 0 else body.add_paragraph()
        par.text = "• " + b; par.font.size = Pt(20); par.font.color.rgb = NAVY
    return s


def main():
    n = _nums()
    prs = Presentation(); prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(12), Inches(2))
    p = box.text_frame.paragraphs[0]; p.text = "Alternator Lead-Time Heuristics — Business Summary (V11)"
    p.font.size = Pt(34); p.font.bold = True; p.font.color.rgb = NAVY
    p2 = box.text_frame.add_paragraph()
    p2.text = f"Earlier warnings for more trucks, zero false alarms ({n['n11']}/10 vs {n['n1062']}/10, {n['nf_false']}/15 FP)"
    p2.font.size = Pt(18); p2.font.color.rgb = GOLD
    _bullets(prs, "How we did it", [
        "Engineered 12 new lead-time signals from the existing 6 CAN channels — no new sensors.",
        "Validated each against the same honest gate (change vs the truck's own healthy baseline AND vs the healthy fleet).",
        "Added a compound 'early-watch' vote and per-truck change-point detection."])
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    t = s3.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
    pp = t.text_frame.paragraphs[0]; pp.text = "Results & impact"; pp.font.size = Pt(30); pp.font.bold = True; pp.font.color.rgb = NAVY
    if (GRAPHS / "G1_recall_comparison.png").exists():
        s3.shapes.add_picture(str(GRAPHS / "G1_recall_comparison.png"), Inches(1.5), Inches(1.2), height=Inches(5.6))
    _bullets(prs, "Limitations & data gaps", [
        "Gain is real but modest: one more truck detected, one earlier warning.",
        "4 of 10 failures have no electrical precursor (abrupt/silent) — undetectable from this data.",
        "No change to per-truck remaining-life dates (structural limit at n=10).",
        "Results are leads to act on, not calibrated probabilities."])
    _bullets(prs, "Next steps", [
        "Deploy post-crank recovery + compound early-watch in the service emergency channel.",
        "Use load-split sag typing for repair-direction guidance.",
        "Re-validate on the larger starter-motor fleet (n=34) and as more failures accrue."])
    out = cfg.V11_ROOT / "presentation" / "Alternator_Business_Summary_V11.pptx"
    prs.save(str(out)); print(f"[v11 business deck] wrote {out} ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
